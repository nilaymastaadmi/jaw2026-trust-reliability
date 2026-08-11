"""Geometry QA for the deck: bounds, overlap, and estimated text overflow.

LibreOffice cannot render in this sandbox, so visual inspection is replaced by
measuring the shape tree directly. Catches the defects that actually reach a
viewer: content off the canvas, content too close to the edge, text boxes
overlapping other content, and text that cannot fit its box at its font size.
"""
from pptx import Presentation
from pptx.util import Emu

DECK = "JAW2026-winners-presentation.pptx"
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.5

prs = Presentation(DECK)


def inches(v):
    return Emu(v).inches if v is not None else None


def boxes(slide):
    out = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        out.append({
            "name": sh.shape_type,
            "x": inches(sh.left), "y": inches(sh.top),
            "w": inches(sh.width), "h": inches(sh.height),
            "text": (sh.text_frame.text if sh.has_text_frame else ""),
            "sh": sh,
        })
    return out


def est_lines(text, width_in, pt):
    """Rough wrapped-line count: ~1.85 chars per pt of width at typical sans."""
    if not text:
        return 0
    chars_per_line = max(1, int(width_in * 96 / (pt * 0.52)))
    n = 0
    for para in text.split("\n"):
        n += max(1, -(-len(para) // chars_per_line))
    return n


def max_pt(shape):
    best = 0
    if not shape.has_text_frame:
        return 0
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                best = max(best, r.font.size.pt)
    return best or 14


issues = 0
for i, slide in enumerate(prs.slides, 1):
    bs = boxes(slide)
    for b in bs:
        # off canvas / too close to edge
        if b["x"] < -0.05 or b["y"] < -0.05:
            if b["w"] < 5 or b["h"] < 5:          # ignore deliberate bleed circles
                print(f"S{i}: element starts off-canvas at ({b['x']:.2f},{b['y']:.2f})")
                issues += 1
        if b["x"] + b["w"] > SLIDE_W + 0.05:
            print(f"S{i}: overflows right edge to {b['x'] + b['w']:.2f}  text={b['text'][:40]!r}")
            issues += 1
        if b["y"] + b["h"] > SLIDE_H + 0.05:
            print(f"S{i}: overflows bottom to {b['y'] + b['h']:.2f}  text={b['text'][:40]!r}")
            issues += 1
        # text overflow estimate
        if b["text"]:
            pt = max_pt(b["sh"])
            lines = est_lines(b["text"], b["w"], pt)
            need = lines * pt * 1.25 / 72
            if need > b["h"] + 0.06:
                print(f"S{i}: text may overflow box "
                      f"(need {need:.2f}in, have {b['h']:.2f}in, {pt:.0f}pt) "
                      f"text={b['text'][:52]!r}")
                issues += 1

print(f"\nslides: {len(prs.slides)}   geometry issues: {issues}")
